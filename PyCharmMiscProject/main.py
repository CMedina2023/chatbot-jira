from fastapi import FastAPI, Form
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
from pptx import Presentation
import os
import json # Necesario para parsear la respuesta JSON de Gemini
import httpx # Para realizar llamadas HTTP asíncronas

app = FastAPI()

# Habilita CORS para permitir solicitudes desde cualquier origen
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

# Función para cargar el contenido de un archivo PowerPoint
def cargar_conocimiento(path):
    texto = ""
    try:
        prs = Presentation(path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texto += shape.text + "\n"
    except Exception as e:
        print(f"❌ Error al cargar el archivo PowerPoint {path}: {e}")
        # Si el archivo no existe o está corrupto, usa un conocimiento de fallback
        return "No se pudo cargar la documentación. Por favor, asegúrate de que 'PLAN de Capacitacion.pptx' existe y es accesible."
    return texto

# Carga la documentación de Jira desde el archivo PowerPoint
conocimiento_jira = cargar_conocimiento("PLAN de Capacitacion.pptx")
# Si el conocimiento_jira está vacío después de cargar, asigna un mensaje por defecto
if not conocimiento_jira.strip():
    conocimiento_jira = "La documentación de Jira no está disponible en este momento."


# Endpoint para consultar al asistente de Jira
@app.post("/api/consultar")
async def consultar_jira(pregunta: str = Form(...)):
    print("📥 Pregunta recibida:", pregunta)

    # Prepara el prompt para el modelo de Gemini
    # El rol 'user' contiene la pregunta y el contexto de conocimiento
    # El rol 'system' para definir el comportamiento del asistente se integra en el 'user' parts aquí para esta llamada
    # ya que gemini-2.0-flash no soporta 'system' role directamente en el historial de chat para este tipo de interacción.
    prompt_parts = [
        {"text": "Eres un asistente de soporte técnico de Jira dentro de una empresa."},
        {"text": "Utiliza la siguiente documentación interna para responder de forma clara:"},
        {"text": conocimiento_jira},
        {"text": f"Pregunta del usuario: {pregunta}"},
        {"text": "Respuesta clara:"}
    ]

    # Prepara el payload para la API de Gemini
    payload = {
        "contents": [
            {
                "role": "user",
                "parts": prompt_parts
            }
        ]
    }

    # URL y clave de API para Gemini. La clave de API es vacía ya que Canvas la inyecta automáticamente.
    api_key = "AIzaSyAk_hIzA0Ts8ul-h14iXrriXTH45K6tjXM" # Canvas inyecta automáticamente la clave de API en tiempo de ejecución.
    api_url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-2.0-flash:generateContent?key={api_key}"

    try:
        # Realiza la llamada a la API de Gemini usando httpx para solicitudes asíncronas
        async with httpx.AsyncClient() as client_http:
            response = await client_http.post(
                api_url,
                headers={"Content-Type": "application/json"},
                json=payload,
                timeout=30.0 # Añade un timeout para evitar que la solicitud se cuelgue
            )
            response.raise_for_status() # Lanza una excepción para errores HTTP (4xx o 5xx)

            result = response.json()
            print("Response from Gemini API:", result)

            # Extrae la respuesta del modelo
            if result.get("candidates") and len(result["candidates"]) > 0 and \
               result["candidates"][0].get("content") and \
               result["candidates"][0]["content"].get("parts") and \
               len(result["candidates"][0]["content"]["parts"]) > 0:
                respuesta_llm = result["candidates"][0]["content"]["parts"][0].get("text", "").strip()
                return {"respuesta": respuesta_llm}
            else:
                print("❌ Estructura de respuesta inesperada de Gemini:", result)
                return {"respuesta": "❌ Error: La API de Gemini no devolvió una respuesta válida."}

    except httpx.RequestError as e:
        print(f"❌ Error de red o conexión al llamar a la API de Gemini: {e}")
        return {"respuesta": f"❌ Error de conexión al servicio de IA: {str(e)}"}
    except httpx.HTTPStatusError as e:
        print(f"❌ Error en la respuesta HTTP de Gemini (Código: {e.response.status_code}): {e.response.text}")
        return {"respuesta": f"❌ Error del servicio de IA: {e.response.status_code} - {e.response.text}"}
    except json.JSONDecodeError as e:
        print(f"❌ Error al decodificar la respuesta JSON de Gemini: {e}")
        return {"respuesta": f"❌ Error al procesar la respuesta del servicio de IA: {str(e)}"}
    except Exception as e:
        print(f"❌ Error inesperado: {e}")
        return {"respuesta": f"❌ Error al generar respuesta: {str(e)}"}

# Monta el directorio 'static' para servir archivos HTML, CSS, JS estáticos
app.mount("/", StaticFiles(directory="static", html=True), name="static")
